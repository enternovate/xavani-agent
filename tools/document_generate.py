#!/usr/bin/env python3

# Copyright (c) 2025-2026 Enternovate.
# MIT License -- See LICENSE file for full terms.
# Built by Enternovate -- Open source. Private. Local.

"""
Document generation tool — pptx / xlsx / docx with quality presets.

Companion to tools/document_tools.py (which reads only). Generation uses
lazy imports so a missing library yields an actionable error string
instead of a crash. Every preset enforces one consistent visual system:
font family, heading scale, accent color, and table styling.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Optional


class _MissingDep(RuntimeError):
    """Raised when an optional writer dependency is not installed."""


# ---------------------------------------------------------------------------
# Style presets: one visual system per document, applied everywhere.
# ---------------------------------------------------------------------------

STYLES: Dict[str, Dict[str, Any]] = {
    "corporate": {
        "font": "Calibri",
        "accent": "1F4E79",
        "heading_scale": {"title": 32, "h1": 24, "h2": 18},
        "table_style": "Medium Style 2 - Accent 1",
    },
    "minimal": {
        "font": "Helvetica",
        "accent": "222222",
        "heading_scale": {"title": 36, "h1": 26, "h2": 20},
        "table_style": "Table Grid",
    },
    "report": {
        "font": "Georgia",
        "accent": "7B2D26",
        "heading_scale": {"title": 30, "h1": 22, "h2": 17},
        "table_style": "Light Grid Accent 1",
    },
}

VALID_KINDS = {"pptx", "xlsx", "docx"}


def _require(module_name: str, extra_hint: str) -> Any:
    try:
        return __import__(module_name)
    except ImportError as exc:
        raise _MissingDep(
            f"'{module_name}' is not installed. Install it with: "
            f"pip install {extra_hint}"
        ) from exc


def _hex(color: str) -> Any:
    from openpyxl.styles import Color

    return Color(rgb=f"FF{color}")


# ---------------------------------------------------------------------------
# pptx
# ---------------------------------------------------------------------------

def _gen_pptx(
    path: Path,
    title: str,
    slides: List[Dict[str, Any]],
    style: Dict[str, Any],
) -> str:
    pptx = _require("pptx", "python-pptx")
    from pptx.util import Pt

    prs = pptx.Presentation()
    font = style["font"]
    h = style["heading_scale"]

    slide_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(slide_layout)
    cover.shapes.title.text = title
    cover.shapes.title.text_frame.paragraphs[0].runs[0].font.name = font
    cover.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(h["title"])

    content_layout = prs.slide_layouts[1]
    for spec in slides:
        slide = prs.slides.add_slide(content_layout)
        head = slide.shapes.title
        head.text = str(spec.get("title", ""))
        run = (
            head.text_frame.paragraphs[0].runs[0]
            if head.text_frame.paragraphs[0].runs else None
        )
        if run is not None:
            run.font.name = font
            run.font.size = Pt(h["h1"])
        body = slide.placeholders[1].text_frame
        lines = [str(x) for x in spec.get("bullets", [])]
        for i, line in enumerate(lines):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = line
            if para.runs:
                para.runs[0].font.name = font
                para.runs[0].font.size = Pt(14)

    prs.save(str(path))
    return f"{len(slides) + 1} slides"


def _gen_xlsx(
    path: Path,
    title: str,
    sheets: List[Dict[str, Any]],
    style: Dict[str, Any],
) -> str:
    openpyxl = _require("openpyxl", "openpyxl")
    from openpyxl.styles import Font  # noqa: F401

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    bold = Font(bold=True, name=style["font"])
    normal = Font(name=style["font"])

    names_used: set = set()
    for i, sheet_spec in enumerate(sheets):
        name = str(sheet_spec.get("name") or f"Sheet{i + 1}")[:31]
        base = name
        suffix = 2
        while name in names_used:
            name = f"{base[:29]}_{suffix}"
            suffix += 1
        names_used.add(name)
        ws = wb.create_sheet(title=name)
        rows = sheet_spec.get("rows", [])
        header = sheet_spec.get("header")
        row_idx = 1
        if header:
            for col, value in enumerate(header, start=1):
                cell = ws.cell(row=row_idx, column=col, value=str(value))
                cell.font = bold
            row_idx += 1
        for row in rows:
            for col, value in enumerate(row, start=1):
                cell = ws.cell(row=row_idx, column=col, value=value)
                cell.font = normal
            row_idx += 1
        ws.freeze_panes = ws.cell(row=(2 if header else 1), column=1)

    wb.save(str(path))
    return f"{len(sheets)} sheet(s)"


def _gen_docx(
    path: Path,
    title: str,
    sections: List[Dict[str, Any]],
    style: Dict[str, Any],
) -> str:
    docx = _require("docx", "python-docx")
    from docx.shared import Pt, RGBColor

    doc = docx.Document()
    font = style["font"]
    h = style["heading_scale"]
    accent = style["accent"]

    def _rgb() -> RGBColor:
        return RGBColor(
            int(accent[0:2], 16), int(accent[2:4], 16), int(accent[4:6], 16)
        )

    heading = doc.add_heading(title, level=0)
    for run in heading.runs:
        run.font.name = font
        run.font.size = Pt(h["title"])
        run.font.color.rgb = _rgb()

    for section in sections:
        head_text = str(section.get("title", ""))
        if head_text:
            h1 = doc.add_heading(head_text, level=1)
            for run in h1.runs:
                run.font.name = font
                run.font.size = Pt(h["h1"])
                run.font.color.rgb = _rgb()
        for paragraph in section.get("paragraphs", []):
            p = doc.add_paragraph(str(paragraph))
            for run in p.runs:
                run.font.name = font
                run.font.size = Pt(11)
        for bullet in section.get("bullets", []):
            item = doc.add_paragraph(str(bullet), style="List Bullet")
            for run in item.runs:
                run.font.name = font
                run.font.size = Pt(11)

    doc.save(str(path))
    return f"{len(sections)} section(s)"


# ---------------------------------------------------------------------------
# Tool entry point
# ---------------------------------------------------------------------------

def generate_document(
    path: str,
    kind: str,
    title: str = "",
    slides: Optional[List[Dict[str, Any]]] = None,
    sheets: Optional[List[Dict[str, Any]]] = None,
    sections: Optional[List[Dict[str, Any]]] = None,
    style: str = "corporate",
) -> Dict[str, Any]:
    """Generate a styled document at path. Returns a result dict."""
    kind = str(kind).strip().lower().lstrip(".")
    if kind not in VALID_KINDS:
        return {
            "ok": False,
            "error": f"Unsupported kind '{kind}'. Use one of: {sorted(VALID_KINDS)}.",
        }
    preset = STYLES.get(str(style).strip().lower())
    if preset is None:
        return {
            "ok": False,
            "error": f"Unknown style '{style}'. Use one of: {sorted(STYLES)}.",
        }
    target = Path(path).expanduser()
    if not target.parent.exists():
        return {"ok": False, "error": f"Directory does not exist: {target.parent}"}
    if target.suffix.lower().lstrip(".") != kind:
        target = target.with_suffix(f".{kind}")

    title = str(title or "Untitled").strip()

    try:
        if kind == "pptx":
            summary = _gen_pptx(target, title, list(slides or []), preset)
        elif kind == "xlsx":
            summary = _gen_xlsx(target, title, list(sheets or []), preset)
        else:
            summary = _gen_docx(target, title, list(sections or []), preset)
    except _MissingDep as exc:
        return {"ok": False, "error": str(exc)}

    return {
        "ok": True,
        "path": str(target),
        "kind": kind,
        "style": str(style),
        "summary": summary,
    }


GENERATE_DOCUMENT_SCHEMA = {
    "name": "generate_document",
    "description": (
        "Generate a styled .pptx, .xlsx, or .docx file. Applies one of the "
        "quality presets (corporate/minimal/report) so fonts, headings, and "
        "colors stay consistent. Provide slides (pptx), sheets (xlsx), or "
        "sections (docx). See the tool schema references."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Output file path."},
            "kind": {"type": "string", "description": "One of: pptx, xlsx, docx."},
            "title": {"type": "string", "description": "Document title."},
            "slides": {
                "type": "array",
                "description": "pptx only: [{title, bullets:[...]}]",
                "items": {"type": "object"},
            },
            "sheets": {
                "type": "array",
                "description": "xlsx only: [{name, header:[...], rows:[[...]]}]",
                "items": {"type": "object"},
            },
            "sections": {
                "type": "array",
                "description": "docx only: [{title, paragraphs:[...], bullets:[...]}]",
                "items": {"type": "object"},
            },
            "style": {
                "type": "string",
                "description": "Preset: corporate (default), minimal, report.",
            },
        },
        "required": ["path", "kind"],
    },
}


def _handle_generate_document(args: Dict[str, Any]) -> str:
    return json.dumps(generate_document(**args), indent=2)


from tools.registry import registry  # noqa: E402

registry.register(
    name="generate_document",
    toolset="files",
    schema=GENERATE_DOCUMENT_SCHEMA,
    handler=_handle_generate_document,
    description="Generate styled pptx/xlsx/docx files.",
    emoji="📝",
)
