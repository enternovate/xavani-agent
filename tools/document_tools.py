# Copyright (c) 2025-2026 Enternovate.
# MIT License -- See LICENSE file for full terms.
# Built by Enternovate -- Open source. Private. Local.

"""Document reading tool (v0.4.0 roadmap U36).

Extracts plain text from common document formats:
  * .txt / .md / .rst / .log / .csv  — dependency-free UTF-8 read
  * .pdf                             — via ``pypdf`` (lazy)
  * .docx                            — via ``python-docx`` (lazy)
  * .xlsx / .xlsm                    — via ``openpyxl`` (lazy)
  * .pptx                            — via ``python-pptx`` (lazy)

Optional parser dependencies are imported lazily and never required at install
time; a missing one yields a clear, actionable error string rather than a crash.
The tool reads only the path it is given and never writes.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict

_PLAINTEXT_EXTS = {".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv", ".json"}


class _MissingDep(RuntimeError):
    """Raised when an optional parser dependency is not installed."""


def _read_pdf(p: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError as exc:  # pragma: no cover - exercised only without the dep
        raise _MissingDep("install 'pypdf' (pip install pypdf) to read PDF files") from exc
    reader = PdfReader(str(p))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(p: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:  # pragma: no cover
        raise _MissingDep("install 'python-docx' (pip install python-docx) to read .docx files") from exc
    doc = Document(str(p))
    return "\n".join(par.text for par in doc.paragraphs)


def _read_xlsx(p: Path) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as exc:  # pragma: no cover
        raise _MissingDep("install 'openpyxl' (pip install openpyxl) to read .xlsx files") from exc
    wb = load_workbook(str(p), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join("" if c is None else str(c) for c in row))
    wb.close()
    return "\n".join(parts)


def _read_pptx(p: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # pragma: no cover
        raise _MissingDep("install 'python-pptx' (pip install python-pptx) to read .pptx files") from exc
    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


_READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".xlsm": _read_xlsx,
    ".pptx": _read_pptx,
}


def read_document(path: str, max_chars: int = 20000) -> Dict[str, Any]:
    """Read ``path`` and return ``{format, chars, text}`` or ``{error}``."""
    p = Path(path).expanduser()
    if not p.is_file():
        return {"error": f"not a file: {path}"}

    ext = p.suffix.lower()
    try:
        if ext in _PLAINTEXT_EXTS:
            text = p.read_text(encoding="utf-8", errors="replace")
        elif ext in _READERS:
            text = _READERS[ext](p)
        else:
            return {"error": f"unsupported document type: {ext or '(none)'}"}
    except _MissingDep as exc:
        return {"error": str(exc)}
    except Exception as exc:  # noqa: BLE001 - surface parser errors to the agent
        return {"error": f"failed to read {ext} document: {exc}"}

    truncated = False
    if max_chars and len(text) > max_chars:
        text = text[:max_chars]
        truncated = True
    result: Dict[str, Any] = {"format": ext.lstrip("."), "chars": len(text), "text": text}
    if truncated:
        result["truncated"] = True
    return result


def _handle_read_document(args: Dict[str, Any]) -> str:
    path = args.get("path", "")
    if not path:
        return json.dumps({"error": "No 'path' provided."})
    max_chars = args.get("max_chars", 20000)
    try:
        max_chars = int(max_chars)
    except (TypeError, ValueError):
        max_chars = 20000
    return json.dumps(read_document(path, max_chars=max_chars), indent=2)


READ_DOCUMENT_SCHEMA: Dict[str, Any] = {
    "name": "read_document",
    "description": (
        "Extract plain text from a document file (.txt/.md/.csv/.json natively; "
        ".pdf/.docx/.xlsx/.pptx via optional parsers). Returns the text plus its "
        "format and character count. Use to read local documents the user references."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "Path to the document file to read."},
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters to return (default 20000; text is truncated beyond this).",
            },
        },
        "required": ["path"],
    },
}


from tools.registry import registry  # noqa: E402

registry.register(
    name="read_document",
    toolset="files",
    schema=READ_DOCUMENT_SCHEMA,
    handler=_handle_read_document,
    description="Read text from documents (pdf/docx/xlsx/pptx/txt/md).",
    emoji="📄",
)
