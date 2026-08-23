#!/usr/bin/env python3

# Copyright (c) 2025-2026 Enternovate.
# MIT License -- See LICENSE file for full terms.

"""Tests for the document generation tool (Task 2.3)."""

import json

import pytest

from tools.document_generate import (
    STYLES,
    generate_document,
)

pptx = pytest.importorskip("pptx")
openpyxl = pytest.importorskip("openpyxl")
docx = pytest.importorskip("docx")


@pytest.fixture
def out(tmp_path):
    return tmp_path / "out"


class TestPptx:
    def test_generates_slides_with_style(self, out):
        result = generate_document(
            path=str(out) + ".pptx",
            kind="pptx",
            title="Q3 Review",
            slides=[{"title": "Revenue", "bullets": ["Up 12%", "Churn down"]}],
            style="corporate",
        )
        assert result["ok"] is True
        assert result["summary"] == "2 slides"
        from pptx import Presentation

        prs = Presentation(result["path"])
        titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
        assert "Q3 Review" in titles[0]
        assert "Revenue" in titles[1]

    def test_suffix_enforced(self, out):
        result = generate_document(
            path=str(out) + ".txt", kind="pptx", title="T"
        )
        assert result["ok"] is True
        assert result["path"].endswith(".pptx")


class TestXlsx:
    def test_generates_sheets_with_header_freeze(self, out):
        result = generate_document(
            path=str(out) + ".xlsx",
            kind="xlsx",
            title="Budget",
            sheets=[{
                "name": "Costs",
                "header": ["Item", "Amount"],
                "rows": [["Hosting", 400], ["Licenses", 900]],
            }],
            style="report",
        )
        assert result["ok"] is True
        wb = openpyxl.load_workbook(result["path"])
        ws = wb["Costs"]
        assert ws.cell(row=1, column=1).value == "Item"
        assert ws.cell(row=2, column=2).value == 400
        assert ws.freeze_panes == "A2"

    def test_duplicate_sheet_names_deduped(self, out):
        result = generate_document(
            path=str(out) + ".xlsx",
            kind="xlsx",
            sheets=[{"name": "Data"}, {"name": "Data"}],
        )
        wb = openpyxl.load_workbook(result["path"])
        assert len(wb.sheetnames) == 2
        assert len(set(wb.sheetnames)) == 2


class TestDocx:
    def test_generates_sections(self, out):
        result = generate_document(
            path=str(out) + ".docx",
            kind="docx",
            title="Policy",
            sections=[
                {"title": "Scope", "paragraphs": ["Applies to all staff."],
                 "bullets": ["Full-time", "Contract"]}
            ],
            style="minimal",
        )
        assert result["ok"] is True
        d = docx.Document(result["path"])
        text = "\n".join(p.text for p in d.paragraphs)
        assert "Scope" in text
        assert "Applies to all staff." in text


class TestValidation:
    def test_bad_kind_rejected(self, out):
        r = generate_document(path=str(out), kind="pdf")
        assert r["ok"] is False and "Unsupported kind" in r["error"]

    def test_bad_style_rejected(self, out):
        r = generate_document(path=str(out) + ".docx", kind="docx", style="gothic")
        assert r["ok"] is False and "Unknown style" in r["error"]

    def test_missing_dir_rejected(self, tmp_path):
        r = generate_document(
            path=str(tmp_path / "nope" / "x.docx"), kind="docx"
        )
        assert r["ok"] is False and "Directory does not exist" in r["error"]

    def test_all_presets_exist(self):
        assert set(STYLES) == {"corporate", "minimal", "report"}
        for preset in STYLES.values():
            assert set(preset["heading_scale"]) == {"title", "h1", "h2"}


class TestRegistryWiring:
    def test_tool_is_registered_and_dispatchable(self, out):
        from tools.registry import registry

        assert "generate_document" in registry.get_all_tool_names()

    def test_handler_returns_json_dict(self, out):
        from tools.document_generate import _handle_generate_document

        raw = _handle_generate_document({
            "path": str(out) + ".docx", "kind": "docx", "title": "H",
        })
        data = json.loads(raw)
        assert data["ok"] is True
